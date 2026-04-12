from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

prs = Presentation()
prs.slide_width = Emu(17475200)
prs.slide_height = Emu(9753600)

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_PURPLE = RGBColor(0x7B, 0x2F, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CARD_BG = RGBColor(0x25, 0x25, 0x45)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, title, content_lines, title_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    shape.shadow.inherit = False
    
    add_text_box(slide, left+0.15, top+0.1, width-0.3, 0.4, title, font_size=16, color=title_color, bold=True)
    y = top + 0.45
    for line in content_lines:
        add_text_box(slide, left+0.15, y, width-0.3, 0.3, line, font_size=11, color=LIGHT_GRAY)
        y += 0.28
    return shape

def add_code_block(slide, left, top, width, height, code_text, title=""):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0D, 0x11, 0x17)
    shape.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    shape.line.width = Pt(1)
    if title:
        add_text_box(slide, left+0.1, top+0.05, width-0.2, 0.3, title, font_size=10, color=ACCENT_BLUE, bold=True)
        add_text_box(slide, left+0.1, top+0.3, width-0.2, height-0.35, code_text, font_size=10, color=GREEN, font_name='Consolas')
    else:
        add_text_box(slide, left+0.1, top+0.05, width-0.2, height-0.1, code_text, font_size=10, color=GREEN, font_name='Consolas')


# =============== SLIDE 1: Title ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
# Decorative bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(14), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, 1, 1.2, 12, 1, '校园二手交易平台', font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.4, 12, 0.8, '业务解读  ·  核心架构  ·  数据表设计  ·  SQL查询分析  ·  扩展方向', font_size=20, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Subtitle line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.3), Inches(5), Inches(0.03))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT_PURPLE; line.line.fill.background()

add_text_box(slide, 1, 3.6, 12, 0.6, '系统设计与技术汇报', font_size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 12, 0.5, '小组成员：XXX  XXX  XXX', font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.0, 12, 0.5, '汇报日期：2026年4月7日', font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# tech badges
badges = ['Vue 3', 'Spring Boot', 'Go / Gin', 'MySQL', 'GORM']
for i, badge in enumerate(badges):
    bx = 3.5 + i * 1.5
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(5.8), Inches(1.3), Inches(0.4))
    s.fill.solid(); s.fill.fore_color.rgb = CARD_BG; s.line.color.rgb = ACCENT_BLUE; s.line.width = Pt(1)
    add_text_box(slide, bx+0.05, 5.82, 1.2, 0.35, badge, font_size=12, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)


# =============== SLIDE 2: Business Analysis ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '业务解读', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(2), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_card(slide, 0.5, 1.2, 4.2, 2.2, '分析对象', [
    '人 → 用户(买家/卖家)：平台的核心参与者。',
    '    需要身份校验(学生认证)保证交易安全。',
    '物 → 商品：系统核心资产。追踪商品状态',
    '    (ON_SALE、SOLD、PENDING_AUDIT)。',
    '事 → 交易履约：绑定买家、卖家、商品的',
    '    业务枢纽，管理订单全生命周期。'
])

add_card(slide, 5.0, 1.2, 4.2, 2.2, '设计理念', [
    '数据一致性：外键约束保证交易关系',
    '业务闭环：从发布→购买→完成的完整生命周期',
    '高扩展性：评论、求购、收藏模块独立解耦',
    '多角色适配：买家、卖家、管理员三位一体'
])

add_card(slide, 9.5, 1.2, 4.2, 2.2, '用户群体', [
    '在校学生(买家)：浏览商品、比价下单',
    '在校学生(卖家)：发布闲置、管理商品',
    '管理员：内容审核、数据监控、平台运营'
])

add_card(slide, 0.5, 3.7, 6.5, 2.0, '系统目标', [
    '交易透明化：精准追踪每件商品"在哪、被谁发布、什么状态"',
    '交易自动化：摒弃线下手动交接，实现发布、下单、评论的系统化管理',
    '决策数据化：通过数据分析热门商品、活跃用户、价格趋势，辅助平台运营'
], title_color=GOLD)

add_card(slide, 7.3, 3.7, 6.5, 2.0, '核心功能', [
    '商品发布与管理 · 智能搜索与分页浏览 · 在线下单与订单追踪',
    '评论互动系统(支持嵌套回复) · 求购信息发布 · 商品收藏',
    '全网天眼雷达比价 · 后台数据仪表盘 · 图片上传与存储'
], title_color=GREEN)


# =============== SLIDE 3: Order Lifecycle ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '业务逻辑：订单与商品状态流转', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(3.5), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

# Order Status Flow
add_text_box(slide, 0.8, 1.3, 5, 0.4, '订单生命周期 (Order Status)', font_size=20, color=ACCENT_BLUE, bold=True)
order_states = [('PENDING\n待支付', ORANGE), ('PAID\n已支付', ACCENT_BLUE), ('SHIPPED\n已发货', ACCENT_PURPLE), ('COMPLETED\n已完成', GREEN)]
for i, (label, col) in enumerate(order_states):
    x = 1.0 + i * 2.8
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(1.8), Inches(0.9))
    s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()
    add_text_box(slide, x+0.1, 1.95, 1.6, 0.8, label, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+1.9), Inches(2.1), Inches(0.7), Inches(0.4))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = LIGHT_GRAY; arrow.line.fill.background()

# Goods Status Flow
add_text_box(slide, 0.8, 3.2, 5, 0.4, '商品状态 (Goods Status)', font_size=20, color=ACCENT_BLUE, bold=True)
goods_states = [('PENDING_AUDIT\n待审核', ORANGE), ('ON_SALE\n上架中', GREEN), ('SOLD\n已售出', ACCENT_PURPLE)]
for i, (label, col) in enumerate(goods_states):
    x = 1.5 + i * 3.5
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.8), Inches(2.0), Inches(0.9))
    s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()
    add_text_box(slide, x+0.1, 3.85, 1.8, 0.8, label, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 2:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+2.1), Inches(4.0), Inches(1.2), Inches(0.4))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = LIGHT_GRAY; arrow.line.fill.background()

add_card(slide, 0.8, 5.2, 12.5, 1.2, '关键流转节点', [
    '触发器：用户发布商品 → 管理员审核通过 → 买家下单 → 卖家发货 → 买家确认收货',
    '数据闭环：Order 状态变为 COMPLETED 标志交易完结，Goods 状态同步更新为 SOLD，形成完整业务循环',
    '买卖分离：通过 user_id(买家) 和 sale_id(卖家) 双外键设计，清晰区分交易双方角色'
], title_color=GOLD)


# =============== SLIDE 4: Core Architecture (ER) ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '系统核心数据架构', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(3), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

# Center: Order table
s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(2.5), Inches(3.5), Inches(1.5))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT_PURPLE; s.line.fill.background()
add_text_box(slide, 5.2, 2.55, 3, 0.4, 'Order (订单表)', font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 5.2, 2.95, 3, 1, 'user_id → 买家\nsale_id → 卖家\norder_no · total · status', font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Left top: User
s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(3), Inches(1.2))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT_BLUE; s.line.fill.background()
add_text_box(slide, 1, 1.55, 2.6, 0.35, 'User (用户表)', font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.9, 2.6, 0.7, 'username(唯一) · password\nname · phone · email · role', font_size=11, color=RGBColor(0x1A,0x1A,0x2E), align=PP_ALIGN.CENTER)

# Left bottom: Purchase
s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.5), Inches(3), Inches(1.2))
s.fill.solid(); s.fill.fore_color.rgb = GREEN; s.line.fill.background()
add_text_box(slide, 1, 3.55, 2.6, 0.35, 'Purchase (求购表)', font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.9, 2.6, 0.7, 'name · address · phone\nuser_id → User', font_size=11, color=RGBColor(0x1A,0x1A,0x2E), align=PP_ALIGN.CENTER)

# Right top: Goods
s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(1.5), Inches(3.2), Inches(1.2))
s.fill.solid(); s.fill.fore_color.rgb = GOLD; s.line.fill.background()
add_text_box(slide, 10, 1.55, 2.8, 0.35, 'Goods (商品表)', font_size=16, color=RGBColor(0x1A,0x1A,0x2E), bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 10, 1.9, 2.8, 0.7, 'name · price · content · status\nuser_id → 卖家(User)', font_size=11, color=RGBColor(0x1A,0x1A,0x2E), align=PP_ALIGN.CENTER)

# Right bottom: Comment  
s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(3.5), Inches(3.2), Inches(1.2))
s.fill.solid(); s.fill.fore_color.rgb = ORANGE; s.line.fill.background()
add_text_box(slide, 10, 3.55, 2.8, 0.35, 'Comment (评论表)', font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 10, 3.9, 2.8, 0.7, 'content · time · pid · root_id\nuser_id → User · fid → Goods', font_size=11, color=RGBColor(0x1A,0x1A,0x2E), align=PP_ALIGN.CENTER)

# Bottom center: Collection
s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(4.8), Inches(3), Inches(1.0))
s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xE9, 0x1E, 0x63); s.line.fill.background()
add_text_box(slide, 5.4, 4.85, 2.6, 0.35, 'Collection (收藏表)', font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 5.4, 5.2, 2.6, 0.5, 'user_id · fid · module', font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, 0.5, 6.2, 13, 0.5, '核心架构特征：以 Order 为引力中心，通过 user_id 和 sale_id 双外键锁定买卖双方关系，确保每笔交易可溯源。', font_size=14, color=LIGHT_GRAY)


# =============== SLIDE 5: Table Structure Detail ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '动态交易表结构设计', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(3), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_card(slide, 0.5, 1.2, 6.5, 3.0, 'Goods (商品表 — 平台核心资产)', [
    '● 业务关联：通过 user_id 外键关联卖家信息',
    '● 状态流转：status 驱动商品上架、审核、售出流程',
    '    ON_SALE(上架) / PENDING_AUDIT(待审) / SOLD(已售)',
    '● 内容管理：content(TEXT) 支持富文本商品描述',
    '● 价格精度：price(DECIMAL 10,2) 保证金额计算准确',
    '● 图片支持：image字段存储图片URL，支持Go上传服务'
], title_color=GOLD)

add_card(slide, 7.3, 1.2, 6.5, 3.0, 'Order (订单表 — 交易枢纽)', [
    '● 双外键设计：user_id(买家) + sale_id(卖家) 锁定交易双方',
    '● 订单唯一性：order_no 唯一约束，防止重复订单',
    '● 状态管理：PENDING → PAID → SHIPPED → COMPLETED',
    '● 金额追踪：total(DECIMAL 10,2) 记录实际成交价',
    '● 冗余优化：good_name 冗余存储商品名，减少联表查询'
], title_color=ACCENT_BLUE)

add_card(slide, 0.5, 4.5, 4.2, 2.0, 'Comment (评论表)', [
    '● 嵌套评论：pid(父评论) + root_id 支持多级回复',
    '● 时间戳：DATETIME DEFAULT CURRENT_TIMESTAMP',
    '● 多模块：fid 关联商品或求购ID'
], title_color=GREEN)

add_card(slide, 5.0, 4.5, 4.2, 2.0, 'Collection (收藏表)', [
    '● 跨模块设计：module字段区分商品/求购',
    '● 用户关联：user_id 绑定收藏用户',
    '● 灵活扩展：fid 支持任意模块收藏'
], title_color=RGBColor(0xE9,0x1E,0x63))

add_card(slide, 9.5, 4.5, 4.2, 2.0, 'Purchase (求购表)', [
    '● 需求发布：用户发布求购需求信息',
    '● 联系方式：name + phone + address',
    '● 用户追溯：user_id 关联发布者'
], title_color=ORANGE)


# =============== SLIDE 6: Table Summary ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '表结构总览', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(2), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

# Table
from pptx.util import Inches as In
table_data = [
    ['表名', '关键字段', '约束说明'],
    ['User', 'id, username, password, name, phone, email, role', 'username 唯一；role 默认 USER'],
    ['Goods', 'id, name, price, content, status, user_id', 'name 非空；price DECIMAL(10,2)；外键关联 User'],
    ['Order', 'id, order_no, good_name, total, status, user_id, sale_id', 'order_no 唯一；双外键(买家+卖家)关联 User'],
    ['Comment', 'id, content, time, user_id, pid, root_id, fid', 'pid/root_id 支持嵌套回复；time 自动时间戳'],
    ['Collection', 'id, user_id, fid, module', 'module 区分收藏类型(商品/求购)'],
    ['Purchase', 'id, name, address, phone, user_id', '外键关联 User；记录求购联系信息'],
]

rows, cols = len(table_data), 3
tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(13), Inches(4.5)).table

for i, row_data in enumerate(table_data):
    for j, cell_text in enumerate(row_data):
        cell = tbl.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.name = 'Microsoft YaHei'
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = LIGHT_GRAY

# Header row styling
for j in range(cols):
    cell = tbl.cell(0, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_PURPLE
# Data rows
for i in range(1, rows):
    for j in range(cols):
        cell = tbl.cell(i, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if i % 2 == 1 else RGBColor(0x1E,0x1E,0x3A)

# Column widths
tbl.columns[0].width = Inches(2)
tbl.columns[1].width = Inches(6)
tbl.columns[2].width = Inches(5)


# =============== SLIDE 7: Single Table Queries ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '单表查询', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(2), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, 0.5, 1.2, 13, 0.4, '查询所有正在上架的商品名称和价格', font_size=16, color=LIGHT_GRAY)
add_code_block(slide, 0.5, 1.6, 6, 0.8, "SELECT name, price FROM goods\nWHERE status = 'ON_SALE';")

add_text_box(slide, 7, 1.2, 6.5, 0.4, '查询价格在50到500之间的商品，按价格升序', font_size=16, color=LIGHT_GRAY)
add_code_block(slide, 7, 1.6, 6.5, 0.8, "SELECT * FROM goods\nWHERE price BETWEEN 50 AND 500\nORDER BY price ASC;")

add_text_box(slide, 0.5, 2.7, 13, 0.4, '查询商品名包含"iPhone"且状态为上架中的商品', font_size=16, color=LIGHT_GRAY)
add_code_block(slide, 0.5, 3.1, 6, 0.8, "SELECT * FROM goods\nWHERE name LIKE '%iPhone%'\nAND status = 'ON_SALE';")

add_text_box(slide, 7, 2.7, 6.5, 0.4, '查询没有填写邮箱的用户，按ID降序', font_size=16, color=LIGHT_GRAY)
add_code_block(slide, 7, 3.1, 6.5, 0.8, "SELECT * FROM user\nWHERE email IS NULL\nORDER BY id DESC;")

add_text_box(slide, 0.5, 4.2, 13, 0.4, '查询共有哪些不同的商品状态', font_size=16, color=LIGHT_GRAY)
add_code_block(slide, 0.5, 4.6, 6, 0.6, "SELECT DISTINCT status FROM goods;")


# =============== SLIDE 8: Group & Aggregate ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '单表查询 | 分组与聚集函数', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(3.5), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_card(slide, 0.5, 1.2, 6.5, 1.0, '哪些卖家的商品数量最多？平均售价如何？', [
], title_color=GOLD)

add_code_block(slide, 0.5, 2.3, 6.5, 2.5, 
    "SELECT\n"
    "    u.name AS seller_name,\n"
    "    COUNT(*) AS goods_count,\n"
    "    ROUND(AVG(g.price), 2) AS avg_price,\n"
    "    MAX(g.price) AS max_price,\n"
    "    MIN(g.price) AS min_price\n"
    "FROM goods g\n"
    "JOIN user u ON g.user_id = u.id\n"
    "GROUP BY g.user_id\n"
    "HAVING COUNT(*) > 5;",
    title="SQL Query"
)

add_card(slide, 7.5, 1.2, 6, 3.5, '分析逻辑', [
    '基于卖家ID进行聚合 (GROUP BY)，',
    '统计每个卖家的：',
    '  · 商品总数 (COUNT)',
    '  · 平均售价 (AVG)',
    '  · 最高价与最低价 (MAX / MIN)',
    '',
    '利用 HAVING 子句过滤出发布超过5件',
    '商品的活跃卖家，用于识别核心卖家。',
    '',
    '通过 JOIN user 表获取卖家真实姓名，',
    '使报表具备可读性。'
], title_color=GREEN)


# =============== SLIDE 9: Join Queries ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '连接查询', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(2), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, 0.5, 1.2, 13, 0.4, '查询所有在售商品及其卖家姓名 (内连接)', font_size=18, color=LIGHT_GRAY, bold=True)
add_code_block(slide, 0.5, 1.7, 13, 1.0,
    "SELECT g.name AS goodName, g.price, g.content, u.name AS sellerName\n"
    "FROM goods g JOIN user u ON g.user_id = u.id\n"
    "WHERE g.status = 'ON_SALE';")

add_text_box(slide, 0.5, 3.0, 13, 0.4, '查询某个买家购买的所有订单及卖家信息 (多表连接)', font_size=18, color=LIGHT_GRAY, bold=True)
add_code_block(slide, 0.5, 3.5, 13, 1.2,
    "SELECT o.order_no, o.good_name, o.total, o.status,\n"
    "       seller.name AS sellerName\n"
    "FROM `order` o\n"
    "JOIN user seller ON o.sale_id = seller.id\n"
    "WHERE o.user_id = :userId;")

add_text_box(slide, 0.5, 5.0, 13, 0.4, '查询某个商品下的所有评论及评论人 (连接查询)', font_size=18, color=LIGHT_GRAY, bold=True)
add_code_block(slide, 0.5, 5.5, 13, 1.0,
    "SELECT c.content, c.time, u.name AS commentator\n"
    "FROM comment c JOIN user u ON c.user_id = u.id\n"
    "WHERE c.fid = :goodsId;")


# =============== SLIDE 10: Multi-table Join ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '连接查询 | 完整订单明细报表', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(3.5), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, 0.5, 1.2, 13, 0.4, '联查订单表、买家、卖家，获取一张完整的订单明细报表', font_size=18, color=LIGHT_GRAY, bold=True)
add_code_block(slide, 0.5, 1.8, 13, 2.5,
    "SELECT\n"
    "    o.order_no       AS '订单号',\n"
    "    o.good_name      AS '商品名称',\n"
    "    o.total          AS '成交金额',\n"
    "    o.status         AS '订单状态',\n"
    "    buyer.name       AS '买家姓名',\n"
    "    buyer.phone      AS '买家电话',\n"
    "    seller.name      AS '卖家姓名',\n"
    "    seller.phone     AS '卖家电话'\n"
    "FROM `order` o\n"
    "JOIN user buyer  ON o.user_id = buyer.id\n"
    "JOIN user seller ON o.sale_id = seller.id\n"
    "ORDER BY o.id DESC;",
    title="完整订单报表 — 三表联查")

add_card(slide, 0.5, 4.6, 13, 1.5, '查询分析', [
    '● 核心思路：Order 表通过 user_id 和 sale_id 两个外键分别关联 User 表，获取买卖双方完整信息',
    '● 别名技巧：同一张 User 表被 JOIN 两次，分别命名为 buyer 和 seller，实现"一表多用"',
    '● 业务价值：一次查询即可生成包含交易双方信息的完整报表，用于客服、对账和运营分析'
], title_color=GOLD)


# =============== SLIDE 11: Nested Queries ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '嵌套查询', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(2), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, 0.5, 1.2, 6.5, 0.4, '查询下过订单的买家详细信息 (IN)', font_size=16, color=LIGHT_GRAY, bold=True)
add_code_block(slide, 0.5, 1.6, 6.5, 1.3,
    "SELECT * FROM user WHERE id IN (\n"
    "    SELECT DISTINCT user_id\n"
    "    FROM `order`\n"
    ");")

add_text_box(slide, 7.5, 1.2, 6, 0.4, '查询价格高于所有已售商品的在售商品', font_size=16, color=LIGHT_GRAY, bold=True)
add_code_block(slide, 7.5, 1.6, 6, 1.3,
    "SELECT name, price FROM goods\n"
    "WHERE status = 'ON_SALE'\n"
    "AND price > ALL (\n"
    "    SELECT price FROM goods\n"
    "    WHERE status = 'SOLD'\n"
    ");")

add_text_box(slide, 0.5, 3.2, 6.5, 0.4, '查询尚未发布过任何商品的用户 (NOT EXISTS)', font_size=16, color=LIGHT_GRAY, bold=True)
add_code_block(slide, 0.5, 3.6, 6.5, 1.3,
    "SELECT username, phone\n"
    "FROM user u\n"
    "WHERE NOT EXISTS (\n"
    "    SELECT 1 FROM goods g\n"
    "    WHERE g.user_id = u.id\n"
    ");")

add_text_box(slide, 7.5, 3.2, 6, 0.4, '查询有评论但没有下过单的活跃用户', font_size=16, color=LIGHT_GRAY, bold=True)
add_code_block(slide, 7.5, 3.6, 6, 1.5,
    "SELECT * FROM user\n"
    "WHERE id IN (\n"
    "    SELECT DISTINCT user_id\n"
    "    FROM comment\n"
    ") AND id NOT IN (\n"
    "    SELECT DISTINCT user_id\n"
    "    FROM `order`\n"
    ");")


# =============== SLIDE 12: Set & Derived Queries ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '集合查询 与 派生表查询', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(3.5), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, 0.5, 1.2, 13, 0.4, '集合查询：将买家和卖家的姓名及联系方式合并，方便群发通知', font_size=16, color=LIGHT_GRAY, bold=True)
add_code_block(slide, 0.5, 1.7, 13, 1.5,
    "SELECT name AS contact_name, phone FROM user\n"
    "WHERE id IN (SELECT DISTINCT user_id FROM `order`)\n"
    "UNION\n"
    "SELECT name AS contact_name, phone FROM user\n"
    "WHERE id IN (SELECT DISTINCT sale_id FROM `order`);",
    title="UNION 集合查询")

add_text_box(slide, 0.5, 3.5, 13, 0.4, '基于派生表的查询：谁是累计消费超过1000元的 VIP 用户？', font_size=18, color=GOLD, bold=True)
add_code_block(slide, 0.5, 4.0, 7, 2.5,
    "SELECT u.username, u.name,\n"
    "       user_stats.total_spend\n"
    "FROM user u\n"
    "JOIN (\n"
    "    SELECT user_id,\n"
    "           SUM(total) AS total_spend\n"
    "    FROM `order`\n"
    "    GROUP BY user_id\n"
    ") AS user_stats\n"
    "ON u.id = user_stats.user_id\n"
    "WHERE user_stats.total_spend > 1000;",
    title="派生表 (Derived Table)")

add_card(slide, 8, 4.0, 5.5, 2.5, '分析逻辑', [
    '通过子查询(派生表)计算每个用户的',
    '总消费金额 SUM(total)，',
    '',
    '在外层查询中筛选出总消费超过',
    '1000 元的 VIP 用户。',
    '',
    '用于精准营销和用户画像分析。'
], title_color=GREEN)


# =============== SLIDE 13: Tech Architecture ===============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 10, 0.6, '下一步架构：系统推进与扩展方向', font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(4), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

# Current architecture
add_card(slide, 0.5, 1.3, 4, 2, '当前架构', [
    '核心关系型数据库 (MySQL)',
    '6张业务表 + 外键约束',
    '存储过程批量数据生成',
    '多类型SQL查询覆盖'
], title_color=ACCENT_BLUE)

# Backend
add_card(slide, 5, 1.3, 4, 2, '业务逻辑后端化', [
    'Java Spring Boot：封装复杂SQL',
    '为独立的 RESTful API 接口',
    'Go Gin：高并发商品分页查询',
    '图片上传、订单管理等服务'
], title_color=GREEN)

# Frontend
add_card(slide, 9.5, 1.3, 4, 2, '操作界面可视化', [
    'Vue 3 + TypeScript 前端',
    '商品瀑布流、求购区、比价雷达',
    '订单管理、收藏夹等用户中心',
    '后台数据仪表盘可视化'
], title_color=GOLD)

# AI upgrade
add_card(slide, 0.5, 3.8, 6.5, 2.5, '智能化升级', [
    '接入大语言模型打造 AI 校园交易助手',
    '全天候解答关于商品信息、交易规则的疑问',
    '根据用户对话意图自动推荐商品',
    '智能内容审核：自动识别违规商品信息',
    '价格预测：基于历史数据预测合理售价'
], title_color=ACCENT_PURPLE)

add_card(slide, 7.5, 3.8, 6, 2.5, '完善数据库', [
    '从底层核心表搭建、多场景业务查询',
    '到接上前端界面与移动端的完整演进',
    '',
    '后续规划：',
    '  · 用户认证(JWT)与权限管理',
    '  · 支付系统对接 · 消息通知',
    '  · 搜索引擎(ES) · 缓存层(Redis)'
], title_color=ORANGE)


# =============== SAVE ===============
output_path = r'd:\Google_projects\校园二手交易平台汇报.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
