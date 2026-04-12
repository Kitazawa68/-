from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r'd:\Google_projects\校园二手交易系统设计.pptx')
print(f"Slide width: {prs.slide_width} EMU = {prs.slide_width/914400:.2f} inches")
print(f"Slide height: {prs.slide_height} EMU = {prs.slide_height/914400:.2f} inches")
print(f"Total slides: {len(prs.slides)}")
print()

for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ({slide.slide_layout.name}) ===")
    for shape in slide.shapes:
        left_in = shape.left / 914400
        top_in = shape.top / 914400
        w_in = shape.width / 914400
        h_in = shape.height / 914400
        print(f"  Shape: {shape.name}, Type: {shape.shape_type}")
        print(f"    Position: left={left_in:.2f}in, top={top_in:.2f}in, w={w_in:.2f}in, h={h_in:.2f}in")
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    font_info = ""
                    if para.runs:
                        r = para.runs[0]
                        font_info = f" [size={r.font.size}, bold={r.font.bold}]"
                    print(f"    TEXT: {text[:100]}{font_info}")
        if shape.has_table:
            table = shape.table
            print(f"    [TABLE {len(table.rows)}x{len(table.columns)}]")
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip()[:30] for cell in row.cells]
                print(f"      Row {row_idx}: {cells}")
    print()
