# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# Use standard 16:9 size matching user's existing PPT
prs.slide_width = Emu(12192000)   # 13.33 inches
prs.slide_height = Emu(6858000)   # 7.50 inches
SW = 13.33
SH = 7.50

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
PURPLE = RGBColor(0x7B, 0x2F, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CARD = RGBColor(0x25, 0x25, 0x45)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
PINK = RGBColor(0xE9, 0x1E, 0x63)

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

def tb(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, al=PP_ALIGN.LEFT, fn='Microsoft YaHei'):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = fn; p.alignment = al
    return box

def mtb(slide, l, t, w, h, lines, sz=14, c=GRAY, fn='Microsoft YaHei'):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = fn
        p.space_after = Pt(4)
    return box

def bar(slide, l, t, w, col=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.04))
    s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()

def card(slide, l, t, w, h, title, lines, tc=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = CARD; s.line.fill.background()
    tb(slide, l+0.15, t+0.12, w-0.3, 0.35, title, sz=15, c=tc, b=True)
    mtb(slide, l+0.15, t+0.5, w-0.3, h-0.6, lines, sz=11, c=GRAY)

def code(slide, l, t, w, h, text, title=""):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x0D,0x11,0x17)
    s.line.color.rgb = RGBColor(0x30,0x36,0x3D); s.line.width = Pt(1)
    y = t + 0.05
    if title:
        tb(slide, l+0.1, y, w-0.2, 0.25, title, sz=10, c=ACCENT, b=True)
        y += 0.25
    tb(slide, l+0.1, y, w-0.2, h-0.3, text, sz=10, c=GREEN, fn='Consolas')

def title_bar(slide, text):
    tb(slide, 0.83, 0.55, 10, 0.55, text, sz=32, c=WHITE, b=True)
    bar(slide, 0.83, 1.15, 2)


# ========== SLIDE 1: Cover ==========
s1 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s1)
bar(s1, 0, 0, SW)
tb(s1, 1.44, 1.8, 10.44, 0.83, '校园二手交易系统', sz=48, c=WHITE, b=True, al=PP_ALIGN.CENTER)
# Connector line
ln = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.28), Inches(2.8), Inches(2.78), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = PURPLE; ln.line.fill.background()
tb(s1, 1.11, 3.05, 11.11, 0.56, '数据库系统设计与技术汇报', sz=20, c=GRAY, al=PP_ALIGN.CENTER)
# Info box
bx = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.21), Inches(4.2), Inches(4.92), Inches(1.2))
bx.fill.solid(); bx.fill.fore_color.rgb = CARD; bx.line.fill.background()
tb(s1, 4.34, 4.35, 4.54, 0.4, '报告人：XXX  XXX  XXX  XXX  XXX', sz=16, c=GRAY, al=PP_ALIGN.CENTER)
tb(s1, 4.34, 4.8, 4.54, 0.4, '2026.4.7', sz=16, c=GRAY, al=PP_ALIGN.CENTER)
# Tech badges
badges = ['Vue 3', 'Spring Boot', 'Go / Gin', 'MySQL']
for i, badge in enumerate(badges):
    bx2 = 3.2 + i * 2.0
    s = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx2), Inches(5.8), Inches(1.6), Inches(0.4))
    s.fill.solid(); s.fill.fore_color.rgb = CARD; s.line.color.rgb = ACCENT; s.line.width = Pt(1)
    tb(s1, bx2+0.05, 5.82, 1.5, 0.35, badge, sz=12, c=ACCENT, al=PP_ALIGN.CENTER)


# ========== SLIDE 2: TOC ==========
s2 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s2)
tb(s2, 0.83, 0.56, 3, 0.56, '目录', sz=36, c=WHITE, b=True)
bar(s2, 0.83, 1.18, 0.83)
toc = ['01  业务解读与系统目标', '02  核心数据架构(ER关系)', '03  数据表结构设计', '04  SQL核心查询展示', '05  系统扩展方向']
for i, item in enumerate(toc):
    y = 1.7 + i * 1.0
    bx = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.83), Inches(y), Inches(11.67), Inches(0.8))
    bx.fill.solid(); bx.fill.fore_color.rgb = CARD; bx.line.fill.background()
    tb(s2, 1.11, y+0.15, 10, 0.5, item, sz=20, c=WHITE, b=True)


# ========== SLIDE 3: Chapter 01 Divider ==========
s3 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s3)
bar(s3, 0, 0, SW)
tb(s3, 3.89, 2.0, 5.56, 0.42, 'CHAPTER 01', sz=18, c=ACCENT, al=PP_ALIGN.CENTER)
tb(s3, 3.74, 2.5, 5.86, 0.83, '业务解读与系统目标', sz=48, c=WHITE, b=True, al=PP_ALIGN.CENTER)
ln = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.97), Inches(3.5), Inches(1.39), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
mtb(s3, 1.81, 3.9, 9.72, 1.11, [
    '阐述校园二手交易系统的分析对象、设计理念、',
    '用户群体与系统核心目标。'
], sz=16, c=GRAY)
for p in s3.shapes[-1].text_frame.paragraphs: p.alignment = PP_ALIGN.CENTER
bar(s3, 0, 7.42, SW)


# ========== SLIDE 4: Business Analysis ==========
s4 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s4)
title_bar(s4, '聚焦痛点，明确系统价值')
card(s4, 0.83, 1.5, 3.61, 3.5, '01 分析对象', [
    '人 → 用户(买家/卖家)：',
    '平台核心参与者，需学生认证。',
    '',
    '物 → 商品：系统核心资产。',
    '追踪状态(上架/已售/待审)。',
    '',
    '事 → 交易履约：绑定买卖双方',
    '的业务枢纽，管理订单全周期。'
], tc=ACCENT)
card(s4, 4.86, 1.5, 3.61, 3.5, '02 设计理念', [
    '数据一致性：',
    '外键约束保证交易关系完整',
    '',
    '业务闭环：',
    '从发布→购买→完成的完整周期',
    '',
    '高扩展性：',
    '评论/求购/收藏模块独立解耦',
], tc=GREEN)
card(s4, 8.89, 1.5, 3.61, 3.5, '03 项目价值', [
    '交易透明化：',
    '精准追踪每件商品状态和归属',
    '',
    '交易自动化：',
    '系统化管理发布、下单、评论',
    '',
    '决策数据化：',
    '分析热门商品、活跃用户趋势',
], tc=GOLD)
# Bottom summary bar
bx = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.83), Inches(5.4), Inches(11.67), Inches(0.83))
bx.fill.solid(); bx.fill.fore_color.rgb = CARD; bx.line.fill.background()
tb(s4, 1.11, 5.55, 11.11, 0.56, '核心目标：通过技术手段解决校园交易信任难题，构建高效、安全的校园二手生态闭环。', sz=16, c=WHITE, b=True, al=PP_ALIGN.CENTER)


# ========== SLIDE 5: User Roles ==========
s5 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s5)
title_bar(s5, '多角色视角下的需求分析')
card(s5, 0.83, 1.5, 3.61, 3.5, '学生（买家）视角', [
    '便捷查找：快速找到物美价廉的闲置',
    '信息透明：了解商品真实状况',
    '安全交易：交易过程有保障',
    '互动咨询：方便与卖家沟通细节'
], tc=ACCENT)
card(s5, 4.86, 1.5, 3.61, 3.5, '卖家视角', [
    '快速发布：便捷上传商品信息',
    '高效曝光：被更多潜在买家看到',
    '交易管理：方便处理订单和收款',
    '信用积累：建立个人信誉'
], tc=GREEN)
card(s5, 8.89, 1.5, 3.61, 3.5, '平台视角', [
    '用户管理：维护信息，确保真实性',
    '商品监管：审核信息，杜绝违规',
    '交易保障：提供支付渠道',
    '数据分析：了解行为，优化功能'
], tc=GOLD)
ln2 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.83), Inches(5.3), Inches(11.67), Inches(0.02))
ln2.fill.solid(); ln2.fill.fore_color.rgb = ACCENT; ln2.line.fill.background()
tb(s5, 0.83, 5.5, 11.67, 0.56, '通过对买家、卖家和平台三方需求的深度洞察，构建安全、高效、透明的校园二手交易生态系统。', sz=14, c=GRAY, al=PP_ALIGN.CENTER)


# ========== SLIDE 6: Order Lifecycle ==========
s6 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s6)
title_bar(s6, '订单与商品状态流转')
# Order Status Flow
tb(s6, 0.83, 1.4, 5, 0.35, '订单生命周期 (Order Status)', sz=16, c=ACCENT, b=True)
ostates = [('PENDING\n待支付', ORANGE), ('PAID\n已支付', ACCENT), ('SHIPPED\n已发货', PURPLE), ('COMPLETED\n已完成', GREEN)]
for i, (label, col) in enumerate(ostates):
    x = 0.83 + i * 2.8
    s = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(2.0), Inches(0.8))
    s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()
    tb(s6, x+0.1, 1.95, 1.8, 0.7, label, sz=13, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    if i < 3:
        ar = s6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+2.1), Inches(2.1), Inches(0.6), Inches(0.35))
        ar.fill.solid(); ar.fill.fore_color.rgb = GRAY; ar.line.fill.background()
# Goods Status Flow
tb(s6, 0.83, 3.0, 5, 0.35, '商品状态 (Goods Status)', sz=16, c=ACCENT, b=True)
gstates = [('PENDING_AUDIT\n待审核', ORANGE), ('ON_SALE\n上架中', GREEN), ('SOLD\n已售出', PURPLE)]
for i, (label, col) in enumerate(gstates):
    x = 1.5 + i * 3.8
    s = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.5), Inches(2.5), Inches(0.8))
    s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()
    tb(s6, x+0.1, 3.55, 2.3, 0.7, label, sz=13, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    if i < 2:
        ar = s6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+2.6), Inches(3.7), Inches(1.1), Inches(0.35))
        ar.fill.solid(); ar.fill.fore_color.rgb = GRAY; ar.line.fill.background()
card(s6, 0.83, 4.7, 11.67, 1.6, '关键流转节点', [
    '触发器：发布商品 → 管理员审核 → 买家下单 → 卖家发货 → 确认收货',
    'Order COMPLETED 标志交易完结，Goods 同步更新为 SOLD',
    '通过 user_id(买家) + sale_id(卖家) 双外键区分交易双方'
], tc=GOLD)


# ========== SLIDE 7: Chapter 02 Divider ==========
s7 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s7)
bar(s7, 0, 0, SW)
tb(s7, 3.89, 2.0, 5.56, 0.42, 'CHAPTER 02', sz=18, c=ACCENT, al=PP_ALIGN.CENTER)
tb(s7, 3.74, 2.5, 5.86, 0.83, '核心数据架构', sz=48, c=WHITE, b=True, al=PP_ALIGN.CENTER)
ln = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.97), Inches(3.5), Inches(1.39), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
mtb(s7, 1.81, 3.9, 9.72, 1.11, [
    '以Order为引力中心，通过外键锁定买卖双方关系，',
    '确保每笔交易可溯源。'
], sz=16, c=GRAY)
for p in s7.shapes[-1].text_frame.paragraphs: p.alignment = PP_ALIGN.CENTER
bar(s7, 0, 7.42, SW)


# ========== SLIDE 8: ER Diagram ==========
s8 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s8)
title_bar(s8, '系统核心数据架构 (ER关系)')

# Center: Order
s = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2.5), Inches(4.33), Inches(1.5))
s.fill.solid(); s.fill.fore_color.rgb = PURPLE; s.line.fill.background()
tb(s8, 4.7, 2.55, 3.9, 0.35, 'Order (订单表)', sz=16, c=WHITE, b=True, al=PP_ALIGN.CENTER)
tb(s8, 4.7, 2.95, 3.9, 0.9, 'user_id → 买家 | sale_id → 卖家\norder_no · total · status', sz=12, c=GRAY, al=PP_ALIGN.CENTER)

# Left: User
s = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.83), Inches(1.5), Inches(3.2), Inches(1.3))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
tb(s8, 1.0, 1.55, 2.8, 0.35, 'User (用户表)', sz=15, c=WHITE, b=True, al=PP_ALIGN.CENTER)
tb(s8, 1.0, 1.9, 2.8, 0.8, 'username(唯一) · password\nname · phone · email · role', sz=11, c=RGBColor(0x1A,0x1A,0x2E), al=PP_ALIGN.CENTER)

# Left bottom: Purchase
s = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.83), Inches(3.5), Inches(3.2), Inches(1.2))
s.fill.solid(); s.fill.fore_color.rgb = GREEN; s.line.fill.background()
tb(s8, 1.0, 3.55, 2.8, 0.35, 'Purchase (求购表)', sz=15, c=WHITE, b=True, al=PP_ALIGN.CENTER)
tb(s8, 1.0, 3.9, 2.8, 0.7, 'name · address · phone\nuser_id → User', sz=11, c=RGBColor(0x1A,0x1A,0x2E), al=PP_ALIGN.CENTER)

# Right: Goods
s = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(1.5), Inches(3.2), Inches(1.3))
s.fill.solid(); s.fill.fore_color.rgb = GOLD; s.line.fill.background()
tb(s8, 9.5, 1.55, 2.8, 0.35, 'Goods (商品表)', sz=15, c=RGBColor(0x1A,0x1A,0x2E), b=True, al=PP_ALIGN.CENTER)
tb(s8, 9.5, 1.9, 2.8, 0.8, 'name · price · content · status\nuser_id → 卖家(User)', sz=11, c=RGBColor(0x1A,0x1A,0x2E), al=PP_ALIGN.CENTER)

# Right bottom: Comment
s = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(3.5), Inches(3.2), Inches(1.2))
s.fill.solid(); s.fill.fore_color.rgb = ORANGE; s.line.fill.background()
tb(s8, 9.5, 3.55, 2.8, 0.35, 'Comment (评论表)', sz=15, c=WHITE, b=True, al=PP_ALIGN.CENTER)
tb(s8, 9.5, 3.9, 2.8, 0.7, 'content · time · pid · root_id\nuser_id→User · fid→Goods', sz=11, c=RGBColor(0x1A,0x1A,0x2E), al=PP_ALIGN.CENTER)

# Bottom: Collection
s = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(4.7), Inches(3.33), Inches(1.0))
s.fill.solid(); s.fill.fore_color.rgb = PINK; s.line.fill.background()
tb(s8, 5.2, 4.75, 2.9, 0.35, 'Collection (收藏表)', sz=15, c=WHITE, b=True, al=PP_ALIGN.CENTER)
tb(s8, 5.2, 5.1, 2.9, 0.5, 'user_id · fid · module', sz=11, c=GRAY, al=PP_ALIGN.CENTER)

tb(s8, 0.83, 6.1, 11.67, 0.5, '核心架构：以 Order 为引力中心，通过 user_id 和 sale_id 双外键锁定买卖双方，确保每笔交易可溯源。', sz=13, c=GRAY, al=PP_ALIGN.CENTER)


# ========== SLIDE 9: Chapter 03 Divider ==========
s9 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s9)
bar(s9, 0, 0, SW)
tb(s9, 3.89, 2.0, 5.56, 0.42, 'CHAPTER 03', sz=18, c=ACCENT, al=PP_ALIGN.CENTER)
tb(s9, 3.0, 2.5, 7.33, 0.83, '数据表结构展示', sz=48, c=WHITE, b=True, al=PP_ALIGN.CENTER)
ln = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.97), Inches(3.5), Inches(1.39), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
mtb(s9, 1.39, 3.9, 10.56, 1.11, [
    '将概念模型转化为物理模型，展示系统核心数据表的具体结构，',
    '包括字段名、数据类型等，呈现数据库的最终实现形式。'
], sz=16, c=GRAY)
for p in s9.shapes[-1].text_frame.paragraphs: p.alignment = PP_ALIGN.CENTER
bar(s9, 0, 7.42, SW)


# ========== SLIDE 10: User Table ==========
s10 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s10)
title_bar(s10, '用户模型表')
tb(s10, 0.83, 1.4, 11.67, 0.42, '用于记录使用系统的普通用户和管理员的核心档案信息，是系统权限管理与身份认证的基础。', sz=16, c=GRAY)
tbl = s10.shapes.add_table(7, 3, Inches(0.83), Inches(2.0), Inches(11.67), Inches(3.5)).table
data = [['字段名','数据类型','作用说明'],
    ['id','int (主键)','用户唯一标识，自增主键'],
    ['username','varchar (唯一)','用户登录账号，具有唯一性'],
    ['password','varchar','用户登录密码，存储加密密文'],
    ['name','varchar','用户真实姓名或昵称'],
    ['phone','varchar','用户手机号码'],
    ['role','varchar','用户角色(USER/ADMIN)，默认USER']]
for i, row in enumerate(data):
    for j, cell_text in enumerate(row):
        cell = tbl.cell(i, j); cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.name = 'Microsoft YaHei'
            p.font.color.rgb = WHITE if i == 0 else GRAY
            p.font.bold = (i == 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE if i == 0 else (CARD if i%2==1 else RGBColor(0x1E,0x1E,0x3A))
tbl.columns[0].width = Inches(2.5); tbl.columns[1].width = Inches(3); tbl.columns[2].width = Inches(6.17)
bx = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.83), Inches(5.8), Inches(11.67), Inches(0.83))
bx.fill.solid(); bx.fill.fore_color.rgb = CARD; bx.line.fill.background()
tb(s10, 1.11, 5.95, 11.11, 0.5, '该表作为系统基础核心表，承载所有用户身份信息，是权限控制、数据隔离和个性化服务的基石。', sz=14, c=GRAY, al=PP_ALIGN.CENTER)


# ========== SLIDE 11: Goods Table ==========
s11 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s11)
title_bar(s11, '商品信息模型表')
bar(s11, 0.83, 1.18, 11.67, RGBColor(0x30,0x30,0x50))
tbl = s11.shapes.add_table(7, 3, Inches(0.97), Inches(1.5), Inches(11.11), Inches(3.2)).table
data = [['字段名','数据类型','作用说明'],
    ['id','int (主键)','商品唯一标识，自增ID'],
    ['name','varchar(100)','商品名称，限制长度确保索引效率'],
    ['price','decimal(10,2)','商品价格，精确存储避免浮点误差'],
    ['content','text','商品详细描述，支持大文本'],
    ['status','varchar(20)','商品状态：ON_SALE/PENDING_AUDIT/SOLD'],
    ['user_id','int (外键)','关联用户表主键，标识商品发布者']]
for i, row in enumerate(data):
    for j, cell_text in enumerate(row):
        cell = tbl.cell(i, j); cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.name = 'Microsoft YaHei'
            p.font.color.rgb = WHITE if i == 0 else GRAY
            p.font.bold = (i == 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE if i == 0 else (CARD if i%2==1 else RGBColor(0x1E,0x1E,0x3A))
tbl.columns[0].width = Inches(2.5); tbl.columns[1].width = Inches(3); tbl.columns[2].width = Inches(5.61)
bx = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.83), Inches(5.0), Inches(11.67), Inches(1.3))
bx.fill.solid(); bx.fill.fore_color.rgb = CARD; bx.line.fill.background()
tb(s11, 1.11, 5.1, 11, 0.35, '核心关联逻辑', sz=16, c=ACCENT, b=True)
tb(s11, 1.11, 5.5, 11, 0.7, '该表记录卖家发布的商品详情，通过user_id与用户表建立一对多关系，实现数据关联查询与权限控制。', sz=14, c=GRAY)


# ========== SLIDE 12: Order Table ==========
s12 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s12)
title_bar(s12, '订单信息模型表')
bar(s12, 0.83, 1.18, 0.83)
tbl = s12.shapes.add_table(8, 3, Inches(0.97), Inches(1.5), Inches(11.11), Inches(3.8)).table
data = [['字段名','数据类型','作用说明'],
    ['id','int (主键)','订单唯一标识，自增ID'],
    ['order_no','varchar(50) 唯一','订单编号，唯一业务单号'],
    ['good_name','varchar(100)','商品名称，冗余存储减少联表'],
    ['total','decimal(10,2)','订单总金额，精确计算'],
    ['user_id','int (外键)','买家ID，关联用户信息表'],
    ['sale_id','int (外键)','卖家ID，关联用户信息表'],
    ['status','varchar(20)','订单状态(PENDING/PAID/SHIPPED/COMPLETED)']]
for i, row in enumerate(data):
    for j, cell_text in enumerate(row):
        cell = tbl.cell(i, j); cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12); p.font.name = 'Microsoft YaHei'
            p.font.color.rgb = WHITE if i == 0 else GRAY
            p.font.bold = (i == 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE if i == 0 else (CARD if i%2==1 else RGBColor(0x1E,0x1E,0x3A))
tbl.columns[0].width = Inches(2.5); tbl.columns[1].width = Inches(3); tbl.columns[2].width = Inches(5.61)
bx = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.83), Inches(5.5), Inches(11.67), Inches(1.3))
bx.fill.solid(); bx.fill.fore_color.rgb = CARD; bx.line.fill.background()
tb(s12, 1.11, 5.6, 11, 0.35, '核心作用：', sz=18, c=ACCENT, b=True)
tb(s12, 1.11, 6.0, 11, 0.7, '该表是交易系统的核心枢纽，通过关联买家(user_id)、卖家(sale_id)和商品信息，完整记录每笔交易。', sz=14, c=GRAY)


# ========== SLIDE 13: Chapter 04 Divider ==========
s13 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s13)
bar(s13, 0, 0, SW)
tb(s13, 3.89, 2.0, 5.56, 0.42, 'CHAPTER 04', sz=18, c=ACCENT, al=PP_ALIGN.CENTER)
tb(s13, 3.0, 2.5, 7.33, 0.83, 'SQL 核心查询展示', sz=48, c=WHITE, b=True, al=PP_ALIGN.CENTER)
ln = s13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.97), Inches(3.5), Inches(1.39), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
mtb(s13, 1.81, 3.9, 9.72, 1.11, [
    '展示单表查询、连接查询、嵌套查询、集合查询与派生表查询，',
    '覆盖系统核心业务场景。'
], sz=16, c=GRAY)
for p in s13.shapes[-1].text_frame.paragraphs: p.alignment = PP_ALIGN.CENTER
bar(s13, 0, 7.42, SW)


# ========== SLIDE 14: Single Table Queries ==========
s14 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s14)
title_bar(s14, '单表查询')
tb(s14, 0.83, 1.4, 5.5, 0.3, '查询所有正在上架的商品名称和价格', sz=13, c=GRAY, b=True)
code(s14, 0.83, 1.75, 5.5, 0.7, "SELECT name, price FROM goods\nWHERE status = 'ON_SALE';")
tb(s14, 6.94, 1.4, 5.56, 0.3, '查询价格在50~500之间的商品', sz=13, c=GRAY, b=True)
code(s14, 6.94, 1.75, 5.56, 0.7, "SELECT * FROM goods\nWHERE price BETWEEN 50 AND 500\nORDER BY price ASC;")
tb(s14, 0.83, 2.7, 5.5, 0.3, '查询商品名包含"iPhone"的在售商品', sz=13, c=GRAY, b=True)
code(s14, 0.83, 3.0, 5.5, 0.7, "SELECT * FROM goods\nWHERE name LIKE '%iPhone%'\nAND status = 'ON_SALE';")
tb(s14, 6.94, 2.7, 5.56, 0.3, '查询没有填写邮箱的用户，按ID降序', sz=13, c=GRAY, b=True)
code(s14, 6.94, 3.0, 5.56, 0.7, "SELECT * FROM user\nWHERE email IS NULL\nORDER BY id DESC;")
tb(s14, 0.83, 4.0, 5.5, 0.3, '查询共有哪些不同的商品状态', sz=13, c=GRAY, b=True)
code(s14, 0.83, 4.3, 5.5, 0.5, "SELECT DISTINCT status FROM goods;")
tb(s14, 6.94, 4.0, 5.56, 0.3, '分组统计每个卖家的商品数量和均价', sz=13, c=GRAY, b=True)
code(s14, 6.94, 4.3, 5.56, 1.5, "SELECT u.name, COUNT(*) AS cnt,\n  ROUND(AVG(g.price),2) AS avg_p,\n  MAX(g.price) AS max_p\nFROM goods g JOIN user u\n  ON g.user_id = u.id\nGROUP BY g.user_id\nHAVING COUNT(*) > 5;", title="GROUP BY + HAVING")


# ========== SLIDE 15: Join Queries ==========
s15 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s15)
title_bar(s15, '连接查询')
tb(s15, 0.83, 1.4, 11.67, 0.3, '查询所有在售商品及其卖家姓名 (内连接)', sz=14, c=GRAY, b=True)
code(s15, 0.83, 1.75, 11.67, 0.9, "SELECT g.name AS goodName, g.price, g.content, u.name AS sellerName\nFROM goods g JOIN user u ON g.user_id = u.id\nWHERE g.status = 'ON_SALE';")
tb(s15, 0.83, 2.9, 11.67, 0.3, '查询某个买家的所有订单及卖家信息 (多表连接)', sz=14, c=GRAY, b=True)
code(s15, 0.83, 3.25, 11.67, 1.0, "SELECT o.order_no, o.good_name, o.total, o.status,\n       seller.name AS sellerName\nFROM `order` o JOIN user seller ON o.sale_id = seller.id\nWHERE o.user_id = :userId;")
tb(s15, 0.83, 4.5, 11.67, 0.3, '查询某个商品下的所有评论及评论人', sz=14, c=GRAY, b=True)
code(s15, 0.83, 4.85, 11.67, 0.9, "SELECT c.content, c.time, u.name AS commentator\nFROM comment c JOIN user u ON c.user_id = u.id\nWHERE c.fid = :goodsId;")
card(s15, 0.83, 6.0, 11.67, 0.7, '以上三条查询已封装为 Spring Boot JPA Repository 方法，通过 RESTful API 对外提供服务。', [], tc=GOLD)


# ========== SLIDE 16: Full Order Report ==========
s16 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s16)
title_bar(s16, '连接查询 | 完整订单明细报表')
tb(s16, 0.83, 1.4, 11.67, 0.3, '联查订单表、买家、卖家，获取完整订单明细报表 (三表联查)', sz=14, c=GOLD, b=True)
code(s16, 0.83, 1.8, 11.67, 2.5,
    "SELECT\n"
    "    o.order_no       AS '订单号',\n"
    "    o.good_name      AS '商品名称',\n"
    "    o.total          AS '成交金额',\n"
    "    o.status         AS '订单状态',\n"
    "    buyer.name       AS '买家姓名',\n"
    "    buyer.phone      AS '买家电话',\n"
    "    seller.name      AS '卖家姓名',\n"
    "    seller.phone     AS '卖家电话'\n"
    "FROM `order` o\n"
    "JOIN user buyer  ON o.user_id = buyer.id\n"
    "JOIN user seller ON o.sale_id = seller.id\nORDER BY o.id DESC;",
    title="完整订单报表 — 同表多次JOIN")
card(s16, 0.83, 4.6, 11.67, 1.6, '查询分析', [
    '● Order 表通过 user_id 和 sale_id 两个外键分别关联 User 表，获取买卖双方完整信息',
    '● 同一张 User 表被 JOIN 两次，分别命名为 buyer 和 seller，实现"一表多用"',
    '● 一次查询即可生成包含交易双方信息的完整报表，用于客服、对账和运营分析'
], tc=GREEN)


# ========== SLIDE 17: Nested Queries ==========
s17 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s17)
title_bar(s17, '嵌套查询')
tb(s17, 0.83, 1.4, 5.5, 0.3, '查询下过订单的买家信息 (IN)', sz=13, c=GRAY, b=True)
code(s17, 0.83, 1.75, 5.5, 1.0, "SELECT * FROM user\nWHERE id IN (\n  SELECT DISTINCT user_id\n  FROM `order`\n);")
tb(s17, 6.94, 1.4, 5.56, 0.3, '查询价格高于所有已售商品的在售品', sz=13, c=GRAY, b=True)
code(s17, 6.94, 1.75, 5.56, 1.0, "SELECT name, price FROM goods\nWHERE status='ON_SALE'\nAND price > ALL (\n  SELECT price FROM goods\n  WHERE status='SOLD'\n);")
tb(s17, 0.83, 3.0, 5.5, 0.3, '查询未发布过商品的用户 (NOT EXISTS)', sz=13, c=GRAY, b=True)
code(s17, 0.83, 3.35, 5.5, 1.0, "SELECT username, phone\nFROM user u\nWHERE NOT EXISTS (\n  SELECT 1 FROM goods g\n  WHERE g.user_id = u.id\n);")
tb(s17, 6.94, 3.0, 5.56, 0.3, '有评论但没有下过单的活跃用户', sz=13, c=GRAY, b=True)
code(s17, 6.94, 3.35, 5.56, 1.2, "SELECT * FROM user\nWHERE id IN (\n  SELECT DISTINCT user_id\n  FROM comment\n) AND id NOT IN (\n  SELECT DISTINCT user_id\n  FROM `order`\n);")


# ========== SLIDE 18: Set & Derived Queries ==========
s18 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s18)
title_bar(s18, '集合查询 与 派生表查询')
tb(s18, 0.83, 1.4, 11.67, 0.3, '集合查询：将买家和卖家联系方式合并，方便群发通知 (UNION)', sz=14, c=GRAY, b=True)
code(s18, 0.83, 1.75, 11.67, 1.2,
    "SELECT name AS contact_name, phone FROM user\nWHERE id IN (SELECT DISTINCT user_id FROM `order`)\nUNION\nSELECT name AS contact_name, phone FROM user\nWHERE id IN (SELECT DISTINCT sale_id FROM `order`);",
    title="UNION 集合查询")
tb(s18, 0.83, 3.2, 11.67, 0.3, '基于派生表查询：谁是累计消费超过1000元的 VIP 用户？', sz=14, c=GOLD, b=True)
code(s18, 0.83, 3.55, 6.5, 2.2,
    "SELECT u.username, u.name,\n"
    "       stats.total_spend\n"
    "FROM user u\n"
    "JOIN (\n"
    "  SELECT user_id,\n"
    "    SUM(total) AS total_spend\n"
    "  FROM `order`\n"
    "  GROUP BY user_id\n"
    ") AS stats\n"
    "ON u.id = stats.user_id\n"
    "WHERE stats.total_spend > 1000;",
    title="派生表 (Derived Table)")
card(s18, 7.7, 3.55, 4.8, 2.2, '分析逻辑', [
    '通过子查询(派生表)计算',
    '每个用户的总消费金额，',
    '',
    '在外层查询中筛选出总消费',
    '超过1000元的VIP用户。',
    '',
    '用于精准营销和用户画像。'
], tc=GREEN)


# ========== SLIDE 19: Chapter 05 Divider ==========
s19 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s19)
bar(s19, 0, 0, SW)
tb(s19, 3.89, 2.0, 5.56, 0.42, 'CHAPTER 05', sz=18, c=ACCENT, al=PP_ALIGN.CENTER)
tb(s19, 3.0, 2.5, 7.33, 0.83, '系统扩展方向', sz=48, c=WHITE, b=True, al=PP_ALIGN.CENTER)
ln = s19.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.97), Inches(3.5), Inches(1.39), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
bar(s19, 0, 7.42, SW)


# ========== SLIDE 20: Future Architecture ==========
s20 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s20)
title_bar(s20, '下一步架构：系统推进与扩展方向')
card(s20, 0.83, 1.5, 3.61, 2.2, '当前架构', [
    '核心关系型数据库(MySQL)',
    '6张业务表 + 外键约束',
    '存储过程批量数据生成',
    '多类型SQL查询覆盖'
], tc=ACCENT)
card(s20, 4.86, 1.5, 3.61, 2.2, '业务逻辑后端化', [
    'Java Spring Boot 封装SQL',
    '为 RESTful API 接口',
    'Go Gin 高并发分页查询',
    '图片上传与订单管理'
], tc=GREEN)
card(s20, 8.89, 1.5, 3.61, 2.2, '操作界面可视化', [
    'Vue 3 + TypeScript 前端',
    '商品瀑布流/求购区/比价',
    '订单管理/收藏夹等用户中心',
    '后台数据仪表盘可视化'
], tc=GOLD)
card(s20, 0.83, 4.1, 5.56, 2.2, '智能化升级', [
    '接入大语言模型打造 AI 校园交易助手',
    '全天候解答商品信息与交易规则疑问',
    '根据对话意图自动推荐商品',
    '智能内容审核 · 价格预测'
], tc=PURPLE)
card(s20, 6.94, 4.1, 5.56, 2.2, '完善数据库与基础设施', [
    '用户认证(JWT) 与 权限管理',
    '支付系统对接 · 消息通知',
    '搜索引擎(ES) · 缓存层(Redis)',
    '移动端适配 · 微服务拆分'
], tc=ORANGE)


# ========== SLIDE 21: Thank You ==========
s21 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s21)
# Decorative circles
c1 = s21.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.39), Inches(-1.39), Inches(5.56), Inches(5.56))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x20,0x20,0x40); c1.line.fill.background()
c2 = s21.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.56), Inches(4.72), Inches(4.17), Inches(4.17))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x20,0x20,0x40); c2.line.fill.background()
tb(s21, 4.0, 2.2, 5.33, 0.83, '谢谢观看', sz=48, c=WHITE, b=True, al=PP_ALIGN.CENTER)
ln = s21.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.97), Inches(3.2), Inches(1.39), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
tb(s21, 0, 3.6, SW, 0.42, '2026.4.7', sz=18, c=GRAY, al=PP_ALIGN.CENTER)
tb(s21, 0, 4.1, SW, 0.42, '汇报人：XXX  XXX  XXX  XXX  XXX', sz=18, c=GRAY, al=PP_ALIGN.CENTER)
tb(s21, 0, 6.5, SW, 0.35, '感谢您的聆听，欢迎批评指正', sz=12, c=RGBColor(0x66,0x66,0x88), al=PP_ALIGN.CENTER)


# ========== SAVE ==========
out = r'd:\Google_projects\校园二手交易平台汇报.pptx'
prs.save(out)
print(f'Done! Saved to: {out}')
print(f'Total slides: {len(prs.slides)}')
