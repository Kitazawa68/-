from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import json

prs = Presentation(r'd:\Google_projects\数据库汽车租赁系统(1).pptx')

print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Total slides: {len(prs.slides)}")
print()

for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ({slide.slide_layout.name}) ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f"  TEXT: {text}")
        if shape.has_table:
            table = shape.table
            print("  [TABLE]")
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                print(f"    Row {row_idx}: {cells}")
        if shape.shape_type == 13:  # Picture
            print(f"  [IMAGE: {shape.name}]")
    print()
