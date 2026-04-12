from pptx import Presentation
prs = Presentation(r'd:\Google_projects\校园二手交易平台汇报.pptx')
print(f"Total slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    title = texts[0] if texts else "(empty)"
    print(f"  Slide {i+1}: {title}")
